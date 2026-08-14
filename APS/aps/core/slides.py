#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""演示引擎：pptx 读写。

模型：slides = [{title, body:[...], layout}]
"""
import os

from pptx import Presentation
from pptx.util import Inches


class SlidesEngine:
    def __init__(self, path: str = None):
        self.path = path
        self.slides = []        # [{title, body, layout}]
        self.active = 0
        self.dirty = False
        if path:
            self._load(path)
        else:
            self._new_default()

    def _new_default(self):
        self.slides = [
            {"title": "APS 演示文稿", "body": ["点击左侧缩略图编辑", "AI 可直接生成/修改本演示文稿"], "layout": "title"},
            {"title": "第二页", "body": ["输入正文内容…"], "layout": "content"},
        ]
        self.dirty = True

    # ---------------- 模型 ----------------
    def add_slide(self, layout="content"):
        self.slides.append({"title": "新幻灯片", "body": [], "layout": layout})
        self.active = len(self.slides) - 1
        self.dirty = True
        return self.active

    def remove_slide(self, idx: int):
        if len(self.slides) > 1:
            del self.slides[idx]
            self.active = min(self.active, len(self.slides) - 1)
            self.dirty = True

    # ---------------- 读写 ----------------
    def _load(self, path: str):
        prs = Presentation(path)
        self.slides = []
        for slide in prs.slides:
            title = ""
            body = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    title = text
                else:
                    body.extend(l for l in text.splitlines() if l.strip())
            self.slides.append({"title": title or "（无标题）", "body": body, "layout": "content"})
        if not self.slides:
            self._new_default()
        self.active = 0
        self.dirty = False

    def save(self, path: str):
        prs = Presentation()
        title_layout = prs.slide_layouts[0]
        content_layout = prs.slide_layouts[1]
        for s in self.slides:
            layout = title_layout if s.get("layout") == "title" else content_layout
            slide = prs.slides.add_slide(layout)
            if slide.shapes.title is not None:
                slide.shapes.title.text = s.get("title", "")
            for i, line in enumerate(s.get("body", [])[:6]):
                body = slide.placeholders[1]
                tf = body.text_frame
                if i == 0:
                    tf.text = line
                else:
                    tf.add_paragraph().text = line
        prs.save(path)
        self.path = path
        self.dirty = False

    def to_text(self) -> str:
        lines = [f"# 演示文稿：{os.path.basename(self.path) if self.path else '新建'}"
                 f"（共 {len(self.slides)} 页）"]
        for i, s in enumerate(self.slides, 1):
            lines.append(f"\n## 第 {i} 页 · {s.get('title', '')}")
            lines.extend(s.get("body", []))
        return "\n".join(lines)
